from projectai.state.state import State
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import re
import matplotlib.pyplot as plt
import seaborn as sns
import pandas as pd
import numpy as np
import io
import sklearn
import statsmodels


def extract_python_code(text):
    """
    Extracts Python code from a markdown-formatted text block.
    """
    match = re.search(r'```python\n(.*?)\n```', text, re.DOTALL)
    return match.group(1).strip() if match else None

def pdfsynthesizer(state:State):
    """
    Extracts visualization code snippets from analysts and executes them.
    """
    figures = []

    print("Completed Analysts:", state["completed_analysts"])

    for viz in state["completed_analysts"]:
        plt.figure(figsize=(10, 6))

        if isinstance(viz, str):
            print(f"Processing string entry as code snippet.")

            # Extract only the Python code
            code = extract_python_code(viz)
            if not code:
                print("No valid Python code found in the text. Skipping.")
                continue  # Skip if there's no valid code

            # Debug extracted code
            print("Extracted Code Snippet:\n", repr(code))
            print("dataframe columns",state['data'].columns)

            # Ensure the code is valid before compilation
            try:
                fig, ax = plt.subplots(figsize=(10, 6))  # Explicitly create a figure & axes
                
                compiled_code = compile(code, "<string>", "exec")
                exec(compiled_code, {
                    'df': state['data'],  
                    'plt': plt, 
                    'sns': sns, 
                    'pd': pd,
                    'np': np,
                    "sklearn": sklearn,
                    "statsmodels": statsmodels  # Attempt to import statsmodels
                })

                fig = plt.gcf()
                figures.append(fig)
            except SyntaxError as e:
                print(f"SyntaxError: {e}")
            except Exception as e:
                print(f"Execution error: {e}")
        else:
            print(f"Warning: Expected a string for viz, got {type(viz)} instead.")
    
    print(figures)
    
    prs = Presentation()
    
    # Add title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Data Analysis Report"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Generated on {pd.Timestamp.now().strftime('%Y-%m-%d')}"
    
    # Add visualization slides
    for i, fig in enumerate(figures):
        buf = io.BytesIO()
        fig.savefig(buf, format='png', dpi=300, bbox_inches='tight')
        buf.seek(0)

        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Visualization {i+1}"

        pic_left = Inches(0.5)
        pic_top = Inches(1.5)
        pic_width = Inches(6)
        slide.shapes.add_picture(buf, pic_left, pic_top, width=pic_width)

        plt.close(fig)
    
    # Add final slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(2.5)
    top = Inches(3)
    width = Inches(5)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You"
    p.alignment = 1
    p.font.size = Pt(40)
    
    prs.save("/home/sridhar/AI_Spring_2025/1. students/s.sunke/automationproject/abc.pptx")
    print(f"Presentation saved as /home/sridhar/AI_Spring_2025/1. students/s.sunke/automationproject/abc.pptx")
    
    return {"figures": figures, "presentation": prs}
