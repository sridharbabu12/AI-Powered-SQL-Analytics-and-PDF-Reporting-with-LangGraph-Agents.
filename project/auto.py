import os
import anthropic
import supabase
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from typing import List, Dict, Any, Optional, Tuple, Union
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import numpy as np
import uuid
from datetime import datetime
import re
from fastapi import HTTPException
from pydantic import BaseModel, ValidationError
from models import ToolUseBlock

class MCPDataAnalysisPresenter:
    def __init__(self, supabase_url: str, supabase_key: str, anthropic_api_key: Optional[str] = None):
        """
        Initialize the MCP-enabled data analysis and presentation system.
        
        Args:
            supabase_url: Your Supabase project URL
            supabase_key: Your Supabase API key
            anthropic_api_key: Anthropic API key (defaults to environment variable)
        """
        # Set up Supabase client
        self.supabase_client = supabase.create_client(supabase_url, supabase_key)
        
        # Set up Anthropic client
        self.client = anthropic.Anthropic(
            api_key=anthropic_api_key or os.environ.get("ANTHROPIC_API_KEY")
        )
        
        # Store table schema information
        self.schema_info = self._get_schema_info()
        
        # Set default visualization styles
        sns.set_theme(style="whitegrid")
        self.color_palette = sns.color_palette("viridis", 10)
        
        # Track the latest analysis for presentation creation
        self.latest_analysis = None
        self.latest_query = None
        self.latest_data = None
        
        # Define MCP tools for database operations and analysis
        self.tools = [
            {
                "name": "query_database",
                "description": "Query the database for information",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "sql_query": {"type": "string", "description": "The SQL query to execute"}
                    },
                    "required": ["sql_query"]
                }
            },
            {
                "name": "get_column_statistics",
                "description": "Calculate statistics for a specific column",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "column_name": {"type": "string", "description": "Column to analyze"},
                        "statistics": {
                            "type": "array", 
                            "items": {"type": "string"},
                            "description": "Statistics to calculate"
                        }
                    },
                    "required": ["column_name", "statistics"]
                }
            },
            {
                "name": "create_visualization",
                "description": "Generate a specific visualization",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "viz_type": {"type": "string", "description": "Type of visualization (bar, line, scatter, histogram, boxplot, heatmap, pie)"},
                        "title": {"type": "string", "description": "Title for the visualization"},
                        "x_column": {"type": "string", "description": "X-axis column"},
                        "y_column": {"type": "string", "description": "Y-axis column (optional for some chart types)"},
                        "hue": {"type": "string", "description": "Column for color grouping (optional)"},
                        "sort_by": {"type": "string", "description": "Column to sort by (optional)"},
                        "ascending": {"type": "boolean", "description": "Sort in ascending order (optional)"},
                        "figsize": {
                            "type": "array", 
                            "items": {"type": "integer"},
                            "description": "Figure size as [width, height] (optional)"
                        }
                    },
                    "required": ["viz_type", "title"]
                }
            }
        ]
    
    def _get_schema_info(self) -> Dict[str, List[Dict[str, Any]]]:
        """Extract schema information from the Supabase database."""
        try:
            # Query to get table names
            tables_query = """
            SELECT table_name 
            FROM information_schema.tables 
            WHERE table_schema = 'public'
            """
            tables_response = self.supabase_client.rpc('execute_sql', {'query': tables_query}).execute()
            
            if not hasattr(tables_response, 'data') or not tables_response.data:
                raise ValueError("No tables found in database")
            
            tables = [row['table_name'] for row in tables_response.data]
            
            schema_info = {}
            for table in tables:
                try:
                    # Query to get column information
                    columns_query = f"""
                    SELECT column_name, data_type, is_nullable 
                    FROM information_schema.columns 
                    WHERE table_schema = 'public' AND table_name = '{table}'
                    """
                    columns_response = self.supabase_client.rpc('execute_sql', {'query': columns_query}).execute()
                    
                    if not hasattr(columns_response, 'data'):
                        print(f"Warning: No column data found for table {table}")
                        continue
                    
                    schema_info[table] = columns_response.data
                    
                    # Get sample data for this table
                    sample_query = f"""
                    SELECT * FROM "{table}" LIMIT 1
                    """
                    sample_response = self.supabase_client.rpc('execute_sql', {'query': sample_query}).execute()
                    if hasattr(sample_response, 'data') and sample_response.data:
                        # Add sample data to schema info
                        for col in schema_info[table]:
                            col_name = col['column_name']
                            if col_name in sample_response.data[0]:
                                col['sample_value'] = sample_response.data[0][col_name]
                                
                except Exception as e:
                    print(f"Error processing table {table}: {str(e)}")
                    continue
                
            if not schema_info:
                raise ValueError("No schema information could be extracted")
            
            return schema_info
            
        except Exception as e:
            raise ValueError(f"Failed to get schema info: {str(e)}")
    
    def query_database(self, user_query: str) -> pd.DataFrame:
        """Execute a natural language query against the database using Claude MCP."""
        try:
            # Store the query for later use
            self.latest_query = user_query
            
            # First, verify the table exists and get column names
            check_query = """
            SELECT column_name 
            FROM information_schema.columns 
            WHERE table_schema = 'public' 
            AND table_name = 'data'
            ORDER BY ordinal_position
            """
            
            columns_df = self._execute_query(check_query)
            print("\nDebug - Available columns:")
            print(columns_df)
            
            if columns_df.empty:
                raise ValueError("Table 'data' does not exist or has no columns")
            
            # Get exact column names from the database
            column_names = columns_df['column_name'].tolist()
            print("\nDebug - Column names from database:")
            for col in column_names:
                print(f"  - {col}")
            
            # Check if table has data with exact column names
            data_check_query = f"""
            SELECT EXISTS (
                SELECT 1 
                FROM data 
                WHERE "{column_names[0]}" IS NOT NULL
                LIMIT 1
            ) as has_data
            """
            
            data_check_df = self._execute_query(data_check_query)
            print("\nDebug - Data check results:")
            print(data_check_df)
            
            if not data_check_df.empty and not data_check_df.iloc[0]['has_data']:
                raise ValueError("No data found in the table")
            
            # Create a detailed schema description with exact column names
            schema_description = self._format_schema_for_prompt()
            
            # Use MCP to generate SQL with more explicit instructions
            message = self.client.messages.create(
                model="claude-3-7-sonnet-20250219",
                max_tokens=1500,
                temperature=0,
                system="""You are an expert PostgreSQL query generator. You must respond with valid JSON containing:
                {
                    "reasoning": "Step-by-step explanation of your approach",
                    "tables_used": ["list", "of", "tables"],
                    "columns_used": ["list", "of", "columns"],
                    "sql_query": "The final SQL query"
                }
                
                Important rules:
                1. Always quote column names with double quotes
                2. Do not include semicolons in SQL queries
                3. Use exact column names from the schema
                4. Start with simple queries and add complexity gradually
                5. Always check if data exists before complex calculations
                """,
                messages=[{
                    "role": "user", 
                    "content": f"""
                    Database Schema:
                    {schema_description}
                    
                    User Query: {user_query}
                    
                    First verify data exists with a simple query, then build complexity.
                    Start with basic SELECT before adding calculations.
                    Always quote column names with double quotes.
                    Do not include semicolons in SQL queries.
                    """
                }]
            )
            
            # Extract and parse the SQL generation response
            response_text = message.content[0].text if message.content else ""
            
            # Debug output
            print("\nDebug - Raw Message Content:")
            print(message.content)
            print("\nDebug - Response Text:")
            print(response_text)
            
            if not response_text:
                raise ValueError("Empty response from Claude")
            
            query_plan = self._parse_mcp_response(response_text)
            sql_query = query_plan["sql_query"].rstrip(';')  # Remove any trailing semicolon
            
            print("\nDebug - Executing SQL query:")
            print(sql_query)
            
            df = self._execute_query(sql_query)
            
            if df.empty:
                # Try a simpler query with exact column names
                simple_query = f"""
                SELECT 
                    "{column_names[0]}" as first_column,
                    "{column_names[1]}" as second_column
                FROM data
                WHERE "{column_names[0]}" IS NOT NULL 
                LIMIT 5
                """
                print("\nDebug - Trying simpler query:")
                print(simple_query)
                df = self._execute_query(simple_query)
                
                if df.empty:
                    raise ValueError("Both complex and simple queries returned no data")
            
            # Store the raw data for later use
            self.latest_data = df
            
            return df
            
        except Exception as e:
            print(f"\nQuery error: {str(e)}")
            if 'message' in locals():
                print("\nFull message content:")
                print(message.model_dump_json(indent=2))
            raise ValueError(f"Query failed: {str(e)}")
    
    def _parse_mcp_response(self, response_text: str) -> Dict[str, Any]:
        """Parse the MCP response to extract the query plan."""
        try:
            # Clean up the response text
            cleaned_text = response_text.strip()
            
            # Check if the response contains JSON
            if "```json" in cleaned_text:
                parts = cleaned_text.split("```json")
                if len(parts) > 1:
                    json_content = parts[1].split("```")[0].strip()
                else:
                    raise ValueError("Malformed JSON block")
            else:
                json_content = cleaned_text
            
            # Clean up the JSON content
            json_content = json_content.replace('\n', ' ').replace('\\n', '\n')
            
            # Escape control characters
            json_content = re.sub(r'[\x00-\x1F\x7F]', '', json_content)  # Remove control characters
            
            # Attempt to parse the JSON content
            try:
                query_plan = json.loads(json_content)
            except json.JSONDecodeError as e:
                print(f"JSON decoding error: {str(e)}")
                print("Raw JSON content:", json_content)  # Print the raw JSON content for debugging
                raise ValueError("Failed to decode JSON content")
            
            # Validate required keys
            required_keys = ["reasoning", "tables_used", "columns_used", "sql_query"]
            missing_keys = [key for key in required_keys if key not in query_plan]
            
            if missing_keys:
                raise ValueError(f"Missing required keys in query plan: {missing_keys}")
            
            return query_plan
            
        except Exception as e:
            print(f"\nError parsing response: {str(e)}")
            print("\nResponse content:")
            print(response_text[:1000])  # Print first 1000 chars for debugging
            raise ValueError(f"Failed to parse MCP response: {str(e)}")
    
    def _format_schema_for_prompt(self) -> str:
        """Format database schema information for inclusion in prompts."""
        schema_text = "Database Schema:\n"
        
        for table, columns in self.schema_info.items():
            schema_text += f"\nTable: {table}\n"
            for col in columns:
                col_name = col['column_name']
                col_type = col['data_type']
                nullable = "NULL" if col['is_nullable'] == 'YES' else "NOT NULL"
                
                schema_text += f"  - Column: \"{col_name}\"\n"
                schema_text += f"    Type: {col_type}\n"
                schema_text += f"    Constraints: {nullable}\n"
                
                if 'sample_value' in col:
                    sample_value = col['sample_value']
                    # Format sample value appropriately
                    if isinstance(sample_value, (str, int, float)):
                        schema_text += f"    Example: {sample_value}\n"
        
        return schema_text
    
    def _execute_query(self, sql_query: str) -> pd.DataFrame:
        """Execute an SQL query against the Supabase database."""
        try:
            # Execute the SQL query using Supabase RPC
            response = self.supabase_client.rpc(
                'execute_sql', 
                {'query': sql_query}
            ).execute()
            
            # Handle Supabase response
            if hasattr(response, 'data'):
                # Extract data from Supabase response
                data = response.data
            else:
                print("No data attribute in response")
                return pd.DataFrame()
            
            # Check if data is None or empty
            if not data:
                print("No data returned from query")
                return pd.DataFrame()
            
            # Create DataFrame from the data
            df = pd.DataFrame(data)
            
            # Process data types appropriately
            for col in df.columns:
                # Handle date columns
                if 'date' in col.lower() or 'time' in col.lower():
                    try:
                        df[col] = pd.to_datetime(df[col], errors='ignore')
                    except Exception as e:
                        print(f"Error converting {col} to datetime: {str(e)}")
                
                # Handle numeric columns with (000s) in name
                elif '(000s)' in col:
                    try:
                        df[col] = pd.to_numeric(df[col], errors='coerce')
                    except Exception as e:
                        print(f"Error converting {col} to numeric: {str(e)}")
                        
            return df
            
        except Exception as e:
            print(f"Error executing query: {str(e)}")
            print(f"Query was: {sql_query}")
            return pd.DataFrame()
    
    def analyze_data(self, df: pd.DataFrame = None, user_query: str = None) -> Dict[str, Any]:
        """
        Analyze the data using MCP with defined tools and structured reasoning.
        
        Args:
            df: Optional DataFrame to analyze, uses latest data if None
            user_query: Optional user query, uses latest query if None
            
        Returns:
            Dictionary with analysis results
        """
        # Use provided data or fall back to stored data
        df = df if df is not None else self.latest_data
        user_query = user_query if user_query is not None else self.latest_query
        
        if df is None or df.empty:
            raise ValueError("No data available for analysis")
        
        # Create a data profile with basic information
        data_profile = {
            "columns": list(df.columns),
            "column_types": {col: str(df[col].dtype) for col in df.columns},
            "row_count": len(df),
            "column_count": len(df.columns),
            "sample_rows": df.head(3).to_dict(orient='records')
        }
        
        # Convert non-serializable values in sample data
        for record in data_profile["sample_rows"]:
            for key, value in record.items():
                if isinstance(value, pd.Timestamp):
                    record[key] = value.strftime('%Y-%m-%d')
        
        # Initial analysis planning through MCP
        print("Starting analysis planning...")
        planning_message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1500,
            temperature=0,
            system="""
            You are an expert data analyst creating a structured analysis plan.
            
            Follow this reasoning process:
            1. UNDERSTAND: Comprehend what insights the user is looking for
            2. PLAN: Create a step-by-step analysis approach
            3. IDENTIFY: Determine which visualizations will best show the patterns
            
            Return a detailed analysis plan in this JSON format:
            {
                "title": "Analysis title based on the query",
                "description": "Brief description of the analysis approach",
                "analysis_steps": [
                    {
                        "step": "Step name",
                        "description": "What this step will reveal",
                        "required_tools": ["tool1", "tool2"]
                    }
                ],
                "visualization_plan": [
                    {
                        "viz_type": "chart type",
                        "title": "Chart title",
                        "description": "What this visualization will show",
                        "key_columns": ["column1", "column2"]
                    }
                ],
                "expected_insights": ["What insights we expect to find"]
            }
            """,
            messages=[{
                "role": "user", 
                "content": f"""
                User Query: {user_query}
                
                Data Profile:
                {json.dumps(data_profile, indent=2)}
                
                Create a structured analysis plan to answer this query.
                """
            }]
        )
        
        # Parse the planning response
        try:
            plan_text = planning_message.content[0].text
            # Handle if response is wrapped in ```json ... ```
            if "```json" in plan_text:
                plan_text = plan_text.split("```json")[1].split("```")[0]
                
            analysis_plan = json.loads(plan_text)
            print("Analysis plan created")
        except Exception as e:
            print(f"Error parsing analysis plan: {str(e)}")
            print(f"Raw plan text: {planning_message.content[0].text}")
            analysis_plan = {
                "title": "Data Analysis",
                "visualization_plan": [],
                "expected_insights": ["Unable to generate analysis plan"]
            }
        
        # Execute the analysis using MCP and tools
        print("Executing analysis with MCP tools...")
        figures = []
        insights = []
        
        # Create tool handling functions
        def handle_get_column_statistics(column_name, statistics=None):
            """Handle calls to the get_column_statistics tool."""
            if statistics is None or len(statistics) == 0:
                statistics = ["mean", "median", "std", "min", "max", "count", "nunique"]
                
            if column_name not in df.columns:
                return {"error": f"Column '{column_name}' not found in dataframe"}
                
            results = {}
            for stat in statistics:
                try:
                    if stat == "mean" and pd.api.types.is_numeric_dtype(df[column_name]):
                        results[stat] = df[column_name].mean()
                    elif stat == "median" and pd.api.types.is_numeric_dtype(df[column_name]):
                        results[stat] = df[column_name].median()
                    elif stat == "std" and pd.api.types.is_numeric_dtype(df[column_name]):
                        results[stat] = df[column_name].std()
                    elif stat == "min":
                        results[stat] = df[column_name].min()
                    elif stat == "max":
                        results[stat] = df[column_name].max()
                    elif stat == "count":
                        results[stat] = df[column_name].count()
                    elif stat == "nunique":
                        results[stat] = df[column_name].nunique()
                except Exception as e:
                    results[stat] = f"Error calculating {stat}: {str(e)}"
                    
            # Convert non-serializable values
            for key, value in results.items():
                if isinstance(value, (np.int64, np.float64)):
                    results[key] = float(value)
                elif isinstance(value, pd.Timestamp):
                    results[key] = value.strftime('%Y-%m-%d')
                    
            return results
        
        def handle_create_visualization(viz_params):
            """Handle calls to the create_visualization tool."""
            try:
                # Extract parameters
                viz_type = viz_params.get("viz_type")
                title = viz_params.get("title", "Visualization")
                x_column = viz_params.get("x_column")
                y_column = viz_params.get("y_column")
                hue = viz_params.get("hue")
                figsize = viz_params.get("figsize", [10, 6])
                
                # Create figure
                plt.figure(figsize=(figsize[0], figsize[1]) if isinstance(figsize, list) else (10, 6))
                
                # Create appropriate visualization
                if viz_type == "bar":
                    if x_column and y_column:
                        sns.barplot(x=df[x_column], y=df[y_column], hue=df[hue] if hue else None)
                    else:
                        return {"error": "Bar charts require x_column and y_column parameters"}
                        
                elif viz_type == "line":
                    if x_column and y_column:
                        sns.lineplot(x=df[x_column], y=df[y_column], hue=df[hue] if hue else None)
                    else:
                        return {"error": "Line charts require x_column and y_column parameters"}
                
                elif viz_type == "scatter":
                    if x_column and y_column:
                        sns.scatterplot(x=df[x_column], y=df[y_column], hue=df[hue] if hue else None)
                    else:
                        return {"error": "Scatter plots require x_column and y_column parameters"}
                
                elif viz_type == "histogram":
                    if x_column:
                        sns.histplot(df[x_column], kde=True)
                    else:
                        return {"error": "Histograms require x_column parameter"}
                
                elif viz_type == "boxplot":
                    if x_column:
                        if y_column:
                            sns.boxplot(x=df[x_column], y=df[y_column], hue=df[hue] if hue else None)
                        else:
                            sns.boxplot(y=df[x_column])
                    else:
                        return {"error": "Box plots require at least x_column parameter"}
                
                elif viz_type == "heatmap":
                    # For heatmaps, we need to create a pivot or correlation matrix
                    if "extra_params" in viz_params and "correlation" in viz_params["extra_params"]:
                        # Create correlation heatmap of numeric columns
                        numeric_df = df.select_dtypes(include=['number'])
                        if numeric_df.empty:
                            return {"error": "No numeric columns available for correlation heatmap"}
                        sns.heatmap(numeric_df.corr(), annot=True, cmap='viridis')
                    else:
                        return {"error": "Heatmaps require specific configuration"}
                
                elif viz_type == "pie":
                    if x_column:
                        # Create value counts for pie chart
                        plt.pie(df[x_column].value_counts(), labels=df[x_column].value_counts().index, autopct='%1.1f%%')
                    else:
                        return {"error": "Pie charts require x_column parameter"}
                
                else:
                    return {"error": f"Unsupported visualization type: {viz_type}"}
                
                # Add title and labels
                plt.title(title)
                if x_column and viz_type != "pie":
                    plt.xlabel(x_column)
                if y_column and viz_type not in ["pie", "histogram"]:
                    plt.ylabel(y_column)
                
                # Adjust layout
                plt.tight_layout()
                
                # Save the figure reference
                fig = plt.gcf()
                figures.append(fig)
                
                return {"status": "success", "message": f"Created {viz_type} visualization: {title}"}
                
            except Exception as e:
                print(f"Error creating visualization: {str(e)}")
                return {"error": f"Visualization error: {str(e)}"}
        
        def handle_filter_data(column, operator, value):
            """Handle calls to the filter_data tool."""
            nonlocal df
            
            try:
                if column not in df.columns:
                    return {"error": f"Column '{column}' not found"}
                
                # Convert value to appropriate type
                col_dtype = df[column].dtype
                if pd.api.types.is_numeric_dtype(col_dtype):
                    try:
                        value = float(value)
                    except ValueError:
                        return {"error": f"Cannot convert '{value}' to numeric for column '{column}'"}
                
                # Apply the filter
                if operator == "==":
                    filtered_df = df[df[column] == value]
                elif operator == "!=":
                    filtered_df = df[df[column] != value]
                elif operator == ">":
                    filtered_df = df[df[column] > value]
                elif operator == "<":
                    filtered_df = df[df[column] < value]
                elif operator == ">=":
                    filtered_df = df[df[column] >= value]
                elif operator == "<=":
                    filtered_df = df[df[column] <= value]
                elif operator == "contains":
                    filtered_df = df[df[column].astype(str).str.contains(str(value))]
                elif operator == "startswith":
                    filtered_df = df[df[column].astype(str).str.startswith(str(value))]
                elif operator == "endswith":
                    filtered_df = df[df[column].astype(str).str.endswith(str(value))]
                else:
                    return {"error": f"Unsupported operator: {operator}"}
                
                # Update the dataframe
                df = filtered_df
                
                return {
                    "status": "success", 
                    "rows_before": len(df), 
                    "rows_after": len(filtered_df),
                    "filter_applied": f"{column} {operator} {value}"
                }
                
            except Exception as e:
                return {"error": f"Filter error: {str(e)}"}
        
        def handle_transform_data(operation, parameters):
            """Handle calls to the transform_data tool."""
            nonlocal df
            
            try:
                if operation == "group_by":
                    # Check required parameters
                    if "by" not in parameters or "aggregations" not in parameters:
                        return {"error": "group_by requires 'by' and 'aggregations' parameters"}
                    
                    by_columns = parameters["by"]
                    aggregations = parameters["aggregations"]
                    
                    # Validate columns
                    for col in by_columns:
                        if col not in df.columns:
                            return {"error": f"Column '{col}' not found for grouping"}
                    
                    for col, agg in aggregations.items():
                        if col not in df.columns:
                            return {"error": f"Column '{col}' not found for aggregation"}
                    
                    # Perform group by
                    df = df.groupby(by_columns).agg(aggregations).reset_index()
                    
                    return {
                        "status": "success",
                        "operation": "group_by",
                        "by_columns": by_columns,
                        "aggregations": aggregations,
                        "result_shape": df.shape
                    }
                
                elif operation == "sort":
                    # Check required parameters
                    if "by" not in parameters:
                        return {"error": "sort requires 'by' parameter"}
                    
                    by_columns = parameters["by"]
                    ascending = parameters.get("ascending", True)
                    
                    # Validate columns
                    if isinstance(by_columns, str):
                        by_columns = [by_columns]
                        
                    for col in by_columns:
                        if col not in df.columns:
                            return {"error": f"Column '{col}' not found for sorting"}
                    
                    # Perform sort
                    df = df.sort_values(by=by_columns, ascending=ascending)
                    
                    return {
                        "status": "success",
                        "operation": "sort",
                        "by_columns": by_columns,
                        "ascending": ascending
                    }
                
                elif operation == "pivot":
                    # Check required parameters
                    if "index" not in parameters or "columns" not in parameters or "values" not in parameters:
                        return {"error": "pivot requires 'index', 'columns', and 'values' parameters"}
                    
                    index = parameters["index"]
                    columns = parameters["columns"]
                    values = parameters["values"]
                    
                    # Validate columns
                    if index not in df.columns:
                        return {"error": f"Column '{index}' not found for pivot index"}
                    if columns not in df.columns:
                        return {"error": f"Column '{columns}' not found for pivot columns"}
                    if values not in df.columns:
                        return {"error": f"Column '{values}' not found for pivot values"}
                    
                    # Perform pivot
                    df = df.pivot(index=index, columns=columns, values=values).reset_index()
                    
                    return {
                        "status": "success",
                        "operation": "pivot",
                        "index": index,
                        "columns": columns,
                        "values": values,
                        "result_shape": df.shape
                    }
                
                elif operation == "resample":
                    # Check required parameters
                    if "date_column" not in parameters or "rule" not in parameters:
                        return {"error": "resample requires 'date_column' and 'rule' parameters"}
                    
                    date_column = parameters["date_column"]
                    rule = parameters["rule"]
                    agg_method = parameters.get("aggregation", "mean")
                    
                    # Validate columns
                    if date_column not in df.columns:
                        return {"error": f"Column '{date_column}' not found for resampling"}
                    
                    # Ensure date column is datetime
                    if not pd.api.types.is_datetime64_dtype(df[date_column]):
                        df[date_column] = pd.to_datetime(df[date_column], errors='coerce')
                    
                    # Set date column as index, resample, and reset index
                    df = df.set_index(date_column)
                    
                    # Apply different aggregation methods based on data types
                    numeric_cols = df.select_dtypes(include=['number']).columns
                    
                    # Handle different aggregation methods
                    if agg_method == "mean":
                        resampled = df[numeric_cols].resample(rule).mean()
                    elif agg_method == "sum":
                        resampled = df[numeric_cols].resample(rule).sum()
                    elif agg_method == "count":
                        resampled = df[numeric_cols].resample(rule).count()
                    else:
                        return {"error": f"Unsupported aggregation method: {agg_method}"}
                    
                    # Reset index to get the date column back
                    df = resampled.reset_index()
                    
                    return {
                        "status": "success",
                        "operation": "resample",
                        "date_column": date_column,
                        "rule": rule,
                        "aggregation": agg_method,
                        "result_shape": df.shape
                    }
                
                else:
                    return {"error": f"Unsupported transform operation: {operation}"}
                
            except Exception as e:
                return {"error": f"Transform error: {str(e)}"}
        
        # Execute the analysis with tool-calling
        analysis_message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=4000,
            temperature=0,
            system="""
            You are an expert data analyst. Your task is to analyze data using tools.
            
            Follow a structured reasoning process:
            1. PLAN: Consider what analysis steps would be most informative
            2. EXPLORE: Use tools to understand the data distribution and relationships
            3. VISUALIZE: Create appropriate visualizations to reveal patterns
            4. INTERPRET: Draw meaningful conclusions from the analysis
            
            Use tools to transform, analyze, and visualize the data. Provide thoughtful 
            analysis based on the results. Work methodically through the analysis plan.
            """,
            messages=[{
                "role": "user",
                "content": f"""
                User Query: {user_query}
                
                Data Profile:
                {json.dumps(data_profile, indent=2)}
                
                # Continuing the analyze_data method where the user's code left off:
                
                Analysis Plan:
                {json.dumps(analysis_plan, indent=2)}
                
                Please analyze this data using the available tools to address the user's query.
                """
            }],
            tools=self.tools
        )
        
        # Debugging: Print the entire analysis message
        print("Analysis Message:", analysis_message)
        
        # Check if tool_use exists in the response
        tool_use_found = False
        for content in analysis_message.content:
            if isinstance(content, ToolUseBlock):
                tool_use_found = True
                tool_use = content.input
                tool_name = content.name
                
                # Execute the appropriate tool
                tool_result = {}
                if tool_name == "query_database":
                    tool_result = self._execute_query(tool_use["sql_query"])
                elif tool_name == "get_column_statistics":
                    tool_result = handle_get_column_statistics(
                        tool_use["column_name"], 
                        tool_use.get("statistics", [])
                    )
                elif tool_name == "create_visualization":
                    tool_result = handle_create_visualization(tool_use)
                elif tool_name == "filter_data":
                    tool_result = handle_filter_data(
                        tool_use["column"], 
                        tool_use["operator"], 
                        tool_use["value"]
                    )
                elif tool_name == "transform_data":
                    tool_result = handle_transform_data(
                        tool_use["operation"], 
                        tool_use["parameters"]
                    )
                
                # Log tool usage and results
                print(f"Tool call: {tool_name}")
                print(f"Parameters: {json.dumps(tool_use, indent=2)}")
                print(f"Result: {json.dumps(str(tool_result)[:100] + '...' if len(str(tool_result)) > 100 else str(tool_result))}")
                
                # Submit tool output back to Claude
                self.client.messages.tool_outputs.create(
                    message_id=analysis_message.id,
                    tool_outputs=[{
                        "tool_call_id": content.id,
                        "output": json.dumps(tool_result)
                    }]
                )
        
        if not tool_use_found:
            print("No tool_use found in the analysis message.")
            # Handle the case where no tool_use is present
            # You can either return a default value or raise a more descriptive error
            raise ValueError("The response does not contain 'tool_use' attribute.")
        
        # Extract key insights from the analysis text
        insights_pattern = r"Key Insights:(.*?)(?:\n\n|\Z)"
        insights_match = re.search(insights_pattern, analysis_text, re.DOTALL)
        if insights_match:
            insights_text = insights_match.group(1).strip()
            # Split insights by bullet points or numbered items
            insights = re.findall(r'(?:^|\n)(?:\d+\.\s*|\*\s*|-\s*)([^\n]+)', insights_text)
            insights = [insight.strip() for insight in insights if insight.strip()]
        
        # Store the analysis results
        self.latest_analysis = {
            "title": analysis_plan.get("title", "Data Analysis"),
            "description": analysis_plan.get("description", ""),
            "insights": insights,
            "figures": figures,
            "raw_analysis": analysis_text
        }
        
        return self.latest_analysis
    
    def create_presentation(self, title: str = None) -> Presentation:
        """
        Create a PowerPoint presentation from the latest analysis.
        
        Args:
            title: Optional presentation title (uses analysis title if None)
            
        Returns:
            PowerPoint presentation object
        """
        if self.latest_analysis is None:
            raise ValueError("No analysis available for presentation creation")
        
        # Create a new presentation
        prs = Presentation()
        
        # Set presentation title
        title = title or self.latest_analysis["title"]
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = self.latest_analysis["description"]
        
        # Add overview slide
        overview_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(overview_slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Key Findings"
        
        # Add insights as bullet points
        tf = content_shape.text_frame
        tf.text = ""
        
        for insight in self.latest_analysis["insights"]:
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
        
        # Add data visualization slides
        for i, fig in enumerate(self.latest_analysis["figures"]):
            # Create a slide for each figure
            slide_layout = prs.slide_layouts[5]  # Layout with title and content
            slide = prs.slides.add_slide(slide_layout)
            
            # Set the title
            title_shape = slide.shapes.title
            title_shape.text = f"Visualization {i+1}"
            
            # Save figure to bytes buffer
            buf = io.BytesIO()
            fig.savefig(buf, format='png')
            buf.seek(0)
            
            # Add the image to the slide
            slide.shapes.add_picture(buf, Inches(1), Inches(2), width=Inches(8))
        
        # Add recommendation slide if insights are available
        if self.latest_analysis["insights"]:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.placeholders[1]
            
            title_shape.text = "Recommendations"
            
            # Create recommendations from insights using MCP
            recommendations_message = self.client.messages.create(
                model="claude-3-7-sonnet-20250219",
                max_tokens=1000,
                temperature=0,
                system="You are a strategic advisor who creates actionable recommendations based on data insights.",
                messages=[{
                    "role": "user",
                    "content": f"""
                    Based on these insights from data analysis:
                    
                    {json.dumps(self.latest_analysis["insights"], indent=2)}
                    
                    Please provide 3-5 actionable recommendations that would help address the business needs.
                    Format each recommendation as a concise bullet point starting with an action verb.
                    """
                }]
            )
            
            recommendations_text = recommendations_message.content[0].text
            
            # Extract bullet points
            recommendations = []
            for line in recommendations_text.split('\n'):
                line = line.strip()
                if line.startswith('•') or line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                    clean_line = re.sub(r'^[•\-*\d\.]+\s*', '', line)
                    if clean_line:
                        recommendations.append(clean_line)
            
            # Add recommendations as bullet points
            tf = content_shape.text_frame
            tf.text = ""
            
            for rec in recommendations:
                p = tf.add_paragraph()
                p.text = rec
                p.level = 0
        
        # Add next steps slide
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title_shape = slide.shapes.title
        content_shape = slide.placeholders[1]
        
        title_shape.text = "Next Steps"
        
        # Generate next steps using MCP
        next_steps_message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1000,
            temperature=0,
            system="You are a project manager creating clear next steps based on data analysis.",
            messages=[{
                "role": "user",
                "content": f"""
                Based on this data analysis related to:
                {self.latest_query}
                
                With insights:
                {json.dumps(self.latest_analysis["insights"], indent=2)}
                
                Please suggest 3-4 concrete next steps for further analysis or action.
                Format each as a clear bullet point with an owner type (e.g., "Data team:", "Business team:").
                """
            }]
        )
        
        next_steps_text = next_steps_message.content[0].text
        
        # Extract bullet points
        next_steps = []
        for line in next_steps_text.split('\n'):
            line = line.strip()
            if line.startswith('•') or line.startswith('-') or line.startswith('*') or re.match(r'^\d+\.', line):
                clean_line = re.sub(r'^[•\-*\d\.]+\s*', '', line)
                if clean_line:
                    next_steps.append(clean_line)
        
        # Add next steps as bullet points
        tf = content_shape.text_frame
        tf.text = ""
        
        for step in next_steps:
            p = tf.add_paragraph()
            p.text = step
            p.level = 0
        
        # Apply a consistent style to all slides
        self._style_presentation(prs)
        
        return prs
    
    def _style_presentation(self, prs: Presentation) -> None:
        """Apply consistent styling to the presentation."""
        # Define colors
        title_color = RGBColor(0, 70, 127)   # Dark blue
        accent_color = RGBColor(0, 112, 192)  # Medium blue
        
        # Style each slide
        for slide in prs.slides:
            # Style title
            if slide.shapes.title:
                title = slide.shapes.title
                title.text_frame.paragraphs[0].font.color.rgb = title_color
                title.text_frame.paragraphs[0].font.bold = True
                title.text_frame.paragraphs[0].font.size = Pt(36)
            
            # Style content
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Skip the title which we've already styled
                        if shape == slide.shapes.title:
                            continue
                            
                        # Style regular content
                        paragraph.font.size = Pt(20)
                        
                        # Style bullet points
                        if paragraph.level > 0:
                            paragraph.font.size = Pt(18)
                            paragraph.font.color.rgb = accent_color
    
    def save_presentation(self, output_path: str) -> str:
        """
        Generate and save a presentation to a file.
        
        Args:
            output_path: Path to save the presentation
            
        Returns:
            Path to the saved presentation
        """
        prs = self.create_presentation()
        
        # If output_path is a directory, create a filename
        if os.path.isdir(output_path):
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"analysis_{timestamp}.pptx"
            output_path = os.path.join(output_path, filename)
        
        # Ensure path has .pptx extension
        if not output_path.endswith('.pptx'):
            output_path += '.pptx'
            
        # Save the presentation
        prs.save(output_path)
        
        print(f"Presentation saved to: {output_path}")
        return output_path
    
    def generate_executive_summary(self) -> str:
        """
        Generate an executive summary of the analysis.
        
        Returns:
            Executive summary as markdown text
        """
        if self.latest_analysis is None or self.latest_query is None:
            raise ValueError("No analysis available for summary generation")
        
        # Generate summary using MCP
        summary_message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1500,
            temperature=0,
            system="""
            You are an expert at creating executive summaries of data analysis results.
            Create a well-structured, concise executive summary in markdown format.
            
            Follow this structure:
            1. Brief introduction with context (1 paragraph)
            2. Key findings section with 3-5 bullet points
            3. Implications for the business (1-2 paragraphs)
            4. Recommended next steps (3-4 bullet points)
            
            Use clear, direct language without technical jargon.
            Be specific about findings and recommendations.
            """,
            messages=[{
                "role": "user",
                "content": f"""
                Original Query: {self.latest_query}
                
                Analysis Title: {self.latest_analysis["title"]}
                
                Key Insights:
                {json.dumps(self.latest_analysis["insights"], indent=2)}
                
                Raw Analysis Text:
                {self.latest_analysis["raw_analysis"][:1500]}
                
                Please generate a professional executive summary in markdown format.
                """
            }]
        )
        
        summary = summary_message.content[0].text
        
        return summary
    
    def export_analysis_to_json(self, output_path: str = None) -> str:
        """
        Export the latest analysis to JSON format.
        
        Args:
            output_path: Optional path to save the JSON file
            
        Returns:
            Path to the saved JSON file or the JSON string if no path provided
        """
        if self.latest_analysis is None:
            raise ValueError("No analysis available for export")
            
        # Create a serializable version of the analysis
        export_data = {
            "query": self.latest_query,
            "title": self.latest_analysis["title"],
            "description": self.latest_analysis["description"],
            "insights": self.latest_analysis["insights"],
            "raw_analysis": self.latest_analysis["raw_analysis"],
            "created_at": datetime.now().isoformat(),
            "data_sample": self.latest_data.head(5).to_dict(orient='records') if self.latest_data is not None else None
        }
        
        # Convert to JSON
        json_data = json.dumps(export_data, indent=2)
        
        # Save to file if path provided
        if output_path:
            # If output_path is a directory, create a filename
            if os.path.isdir(output_path):
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"analysis_{timestamp}.json"
                output_path = os.path.join(output_path, filename)
            
            # Ensure path has .json extension
            if not output_path.endswith('.json'):
                output_path += '.json'
                
            # Save the JSON file
            with open(output_path, 'w') as f:
                f.write(json_data)
                
            print(f"Analysis exported to: {output_path}")
            return output_path
        
        # Return JSON string if no path provided
        return json_data
    
    def interactive_analysis(self, user_query: str) -> Dict[str, Any]:
        """
        Perform end-to-end interactive analysis from query to insights.
        
        Args:
            user_query: Natural language query for analysis
            
        Returns:
            Dictionary with analysis results and visualization references
        """
        print(f"\nBeginning interactive analysis for: {user_query}")
        
        # Step 1: Query the database
        print("\nQuerying database...")
        df = self.query_database(user_query)
        if df.empty:
            print("No data returned from query")
            return {"error": "No data returned from database query"}
            
        print(f"Retrieved {len(df)} rows with {len(df.columns)} columns")
        
        # Step 2: Analyze the data
        print("\nAnalyzing data...")
        analysis = self.analyze_data(df, user_query)
        
        # Step 3: Generate a summary
        print("\nGenerating executive summary...")
        summary = self.generate_executive_summary()
        
        # Return the complete analysis package
        return {
            "query": user_query,
            "data": df,
            "analysis": analysis,
            "summary": summary,
            "visualizations": [fig for fig in analysis.get("figures", [])]
        }


# Example usage
if __name__ == "__main__":
    # Set up environment variables for API keys
    supabase_url = os.environ.get("SUPABASE_URL")
    supabase_key = os.environ.get("SUPABASE_KEY")
    
    # Create the analyzer
    analyzer = MCPDataAnalysisPresenter(supabase_url, supabase_key)
    
    # Run an interactive analysis
    results = analyzer.interactive_analysis(
        "What insights can be drawn from the National AA (000s) and Coverage AA (000s) data in terms of audience engagement and growth over time?"
    )
    
    # Generate and save a presentation
    analyzer.save_presentation("./output")
    
    # Print the executive summary
    print("\n" + "="*50)
    print("EXECUTIVE SUMMARY")
    print("="*50)
    print(results["summary"])