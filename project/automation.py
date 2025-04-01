import os
import anthropic
import supabase
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from typing import List, Dict, Any, Optional
import io
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import numpy as np

class DataAnalysisPresenter:
    def __init__(self, supabase_url: str, supabase_key: str, anthropic_api_key: Optional[str] = None):
        """
        Initialize the data analysis and presentation system with Anthropic MCP.
        
        Args:
            supabase_url: Your Supabase project URL
            supabase_key: Your Supabase API key
            anthropic_api_key: Anthropic API key (defaults to environment variable)
        """
        # Set up Supabase client
        self.supabase_client = supabase.create_client(supabase_url, supabase_key)
        
        # Set up Anthropic client
        self.client = anthropic.Anthropic(
            api_key=anthropic_api_key or os.environ.get("ANTHROPIC_API_KEY")
        )
        
        # Store table schema information
        self.schema_info = self._get_schema_info()
        
        # Set default visualization styles
        sns.set_theme(style="whitegrid")
        self.color_palette = sns.color_palette("viridis", 10)
        
        # Track the latest analysis for presentation creation
        self.latest_analysis = None
        self.latest_query = None
        self.latest_data = None
    
    def _get_schema_info(self) -> Dict[str, List[Dict[str, Any]]]:
        """Extract schema information from the Supabase database."""
        # Query to get table names
        tables_query = """
        SELECT table_name 
        FROM information_schema.tables 
        WHERE table_schema = 'public'
        """
        tables_response = self.supabase_client.rpc('execute_sql', {'query': tables_query}).execute()
        
        if hasattr(tables_response, 'error') and tables_response.error:
            raise Exception(f"Error fetching tables: {tables_response.error}")
        
        tables = [row['table_name'] for row in tables_response.data]
        
        schema_info = {}
        for table in tables:
            # Query to get column information
            columns_query = f"""
            SELECT column_name, data_type, is_nullable 
            FROM information_schema.columns 
            WHERE table_schema = 'public' AND table_name = '{table}'
            """
            columns_response = self.supabase_client.rpc('execute_sql', {'query': columns_query}).execute()
            
            if hasattr(columns_response, 'error') and columns_response.error:
                raise Exception(f"Error fetching columns for {table}: {columns_response.error}")
            
            schema_info[table] = columns_response.data
            
        return schema_info
    
    def _generate_sql_from_query(self, user_query: str) -> str:
        """Generate SQL from a natural language query using Claude."""
        # Create a detailed schema description with example values
        schema_description = "Database Schema:\n"
        for table, columns in self.schema_info.items():
            schema_description += f"\nTable: {table}\n"
            # Get sample data for this table
            sample_query = f"""
            SELECT * FROM "{table}" LIMIT 1
            """
            try:
                sample_response = self.supabase_client.rpc('execute_sql', {'query': sample_query}).execute()
                sample_data = sample_response.data[0] if sample_response.data else {}
            except Exception:
                sample_data = {}
            
            for col in columns:
                col_name = col['column_name']
                col_type = col['data_type']
                nullable = "NULL" if col['is_nullable'] == 'YES' else "NOT NULL"
                sample_value = sample_data.get(col_name, '')
                schema_description += f"  - Column: \"{col_name}\"\n"
                schema_description += f"    Type: {col_type}\n"
                schema_description += f"    Constraints: {nullable}\n"
                if sample_value:
                    schema_description += f"    Example: {sample_value}\n"
        
        # Send the prompt to Claude
        message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1000,
            temperature=0,
            system="""
            You are an expert PostgreSQL query generator. Generate valid PostgreSQL queries for a Supabase database.
            
            Important rules:
            1. ONLY use columns that exist in the schema
            2. Use exact column names as shown in the schema (case-sensitive)
            3. Always enclose table and column names in double quotes when they:
               - Contain spaces
               - Contain special characters
               - Use mixed case
               - Are PostgreSQL keywords
            4. Do NOT include semicolons at the end of queries
            5. Do NOT prefix column names with table names unless joining multiple tables
            6. Include WHERE clauses to filter data as needed
            7. Return ONLY the SQL query without any explanations
            
            Before generating the query:
            - Verify all referenced columns exist in the schema
            - Check column data types for appropriate comparisons
            - Use proper table aliases if needed
            - Ensure proper quoting of identifiers
            """,
            messages=[
                {"role": "user", "content": f"""
                {schema_description}
                
                Generate a PostgreSQL query (without semicolon) for this request:
                {user_query}
                
                Remember to:
                1. Only use columns that exist in the schema above
                2. Use exact column names with proper case
                3. Enclose column names containing spaces or special characters in double quotes
                4. Do NOT prefix columns with table names unless necessary for joins
                """}
            ]
        )
        
        # Extract and clean the SQL query
        sql_query = message.content[0].text.strip()
        
        # Remove any trailing semicolons
        sql_query = sql_query.rstrip(';')
        
        # Log the generated query for debugging
        print(f"Generated SQL query: {sql_query}")
        
        return sql_query
    
    def query_database(self, user_query: str) -> pd.DataFrame:
        """
        Execute a natural language query against the database and return results as a DataFrame.
        
        Args:
            user_query: The user's natural language query
            
        Returns:
            pandas DataFrame with the query results
        """
        # Store the query for later use
        self.latest_query = user_query
        
        # Generate SQL from the user query
        sql_query = self._generate_sql_from_query(user_query)
        
        try:
            # Execute the SQL query using Supabase RPC
            response = self.supabase_client.rpc(
                'execute_sql', 
                {'query': sql_query}
            ).execute()
            
            # Debug logging
            #print("Response type:", type(response))
            #print("Response content:", response)
            
            # Handle Supabase response
            if hasattr(response, 'data'):
                # Extract data from Supabase response
                data = response.data
            else:
                print("No data attribute in response")
                return pd.DataFrame()
            
            # Check if data is None or empty
            if not data:
                print("No data returned from query")
                return pd.DataFrame()
            
            try:
                # Create DataFrame from the data
                df = pd.DataFrame(data)
                
                # Convert date columns to datetime
                date_columns = [col for col in df.columns if 'date' in col.lower()]
                for col in date_columns:
                    try:
                        df[col] = pd.to_datetime(df[col])
                    except Exception as e:
                        print(f"Error converting {col} to datetime: {str(e)}")
                
                # Convert numeric columns
                numeric_columns = [col for col in df.columns if '(000s)' in col]
                for col in numeric_columns:
                    try:
                        df[col] = pd.to_numeric(df[col], errors='coerce')
                    except Exception as e:
                        print(f"Error converting {col} to numeric: {str(e)}")
                
                # Store the raw data for later use
                self.latest_data = df
                
                return df
                
            except Exception as e:
                print(f"Error creating DataFrame: {str(e)}")
                print(f"Data type: {type(data)}")
                print(f"Data content: {data}")
                return pd.DataFrame()
            
        except Exception as e:
            print(f"Error executing query: {str(e)}")
            print(f"Generated SQL query was: {sql_query}")
            return pd.DataFrame()
    
    def _determine_analysis_type(self, df: pd.DataFrame, user_query: str) -> Dict[str, Any]:
        """
        Use LLM to determine appropriate analysis approach based on the data and query.
        
        Args:
            df: The pandas DataFrame containing the query results
            user_query: The original user query
            
        Returns:
            Dictionary with analysis plan
        """
        # Create a data profile
        data_profile = {
            "columns": list(df.columns),
            "column_types": {col: str(df[col].dtype) for col in df.columns},
            "row_count": len(df),
            "sample_data": df.head(5).to_dict(orient='records'),
            "summary_stats": df.describe().to_dict()
        }

        message = self.client.messages.create(
            model="claude-3-7-sonnet-20250219",
            max_tokens=1500,
            temperature=0,
            system="""
            You are an expert data analyst. Given a dataset and user query, determine the most appropriate 
            analyses to perform. Return a JSON object with:
            
            {
                "title": "Analysis title based on the query",
                "analysis_approach": [
                    {
                        "type": "analysis type",
                        "description": "What this analysis will reveal",
                        "required_columns": ["columns needed"],
                        "statistical_methods": ["methods to use"]
                    }
                ],
                "expected_insights": ["What insights we expect to find"],
                "visualization_types": ["Recommended visualization types"]
            }
            
            Return only valid JSON.
            """,
            messages=[
                {"role": "user", "content": f"""
                User Query: {user_query}
                
                Dataset Profile:
                {json.dumps(data_profile, indent=2)}
                
                What analyses would best answer this query?
                """}
            ]
        )

        return json.loads(message.content[0].text)
    
    def analyze_data(self, df: pd.DataFrame = None, user_query: str = None) -> Dict[str, Any]:
        """
        Analyze the data based on the user query and generate visualizations using LLM.
        """
        # Use provided data or fall back to stored data
        df = df if df is not None else self.latest_data
        user_query = user_query if user_query is not None else self.latest_query
        
        if df is None or df.empty:
            raise ValueError("No data available for analysis")

        # Create a data profile with JSON-serializable values
        data_profile = {
            "columns": list(df.columns),
            "column_types": {col: str(df[col].dtype) for col in df.columns},
            "row_count": len(df),
            "sample_data": df.head(5).to_dict(orient='records'),
        }
        
        # Convert datetime objects to strings in sample data
        for record in data_profile["sample_data"]:
            for key, value in record.items():
                if isinstance(value, pd.Timestamp):
                    record[key] = value.strftime('%Y-%m-%d')

        try:
            # Get visualization code from Claude
            message = self.client.messages.create(
                model="claude-3-7-sonnet-20250219",
                max_tokens=4000,
                temperature=0,
                system="""
                You are an expert data visualization developer. Generate Python code using matplotlib and seaborn
                to create insightful visualizations based on the data and user query.
                
                The visualizations should:
                1. Use clear, professional styling with proper labels and titles
                2. Include appropriate color schemes and legends
                3. Handle data types correctly (dates, numbers, categories)
                4. Use the most suitable chart types for the analysis
                5. Show meaningful insights about the data
                
                Return a JSON object with:
                {
                    "title": "Analysis title",
                    "visualizations": [
                        {
                            "description": "Description of what this visualization shows",
                            "code": "Python code to create the visualization",
                            "analysis": [
                                "Key observation 1 about this visualization",
                                "Key observation 2 about this visualization",
                                "Interpretation of trends or patterns shown"
                            ]
                        }
                    ],
                    "insights": [
                        "Key insight 1",
                        "Key insight 2"
                    ]
                }
                
                For each visualization, include 2-3 specific observations about:
                - Trends or patterns shown in the data
                - Notable outliers or anomalies
                - Relationships between variables
                - Business implications of the findings
                """,
                messages=[{
                    "role": "user", 
                    "content": f"""
                    User Query: {user_query}
                    Data Profile: {json.dumps(data_profile, indent=2)}
                    
                    Generate visualization code to analyze this data, focusing on:
                    - Trends in National and Coverage AA metrics
                    - Year-over-year and weekly growth patterns
                    - Performance comparisons across different time periods
                    - Key relationships between metrics
                    """
                }]
            )

            # Parse Claude's response
            try:
                # Extract JSON from the response
                json_str = message.content[0].text
                if "```json" in json_str:
                    json_str = json_str.split("```json")[-1].split("```")[0]
                
                analysis_plan = json.loads(json_str)
                self.latest_analysis = analysis_plan
                
            except Exception as e:
                print(f"Error parsing Claude's response: {str(e)}")
                print("Claude's response:", message.content[0].text)
                return {
                    "title": "Analysis Error",
                    "figures": [],
                    "insights": ["Error parsing analysis response"],
                    "analysis_plan": {},
                    "data": df
                }

            figures = []
            
            # Execute each visualization code snippet
            for viz in analysis_plan.get("visualizations", []):
                plt.figure(figsize=(10, 6))
                try:
                    # Clean up the code string
                    code = viz["code"].strip()
                    if code.startswith('```python'):
                        code = code.split('```python')[1]
                    if code.endswith('```'):
                        code = code[:-3]
                    
                    # Execute the visualization code
                    exec(code, {
                        'df': df, 
                        'plt': plt, 
                        'sns': sns, 
                        'pd': pd,
                        'np': np
                    })
                    
                    # Capture the current figure
                    fig = plt.gcf()
                    figures.append(fig)
                    
                except Exception as e:
                    print(f"Error creating visualization: {str(e)}")
                    plt.close()
                    continue

            return {
                "title": analysis_plan.get("title", "Data Analysis"),
                "figures": figures,
                "insights": analysis_plan.get("insights", []),
                "analysis_plan": analysis_plan,
                "data": df
            }
            
        except Exception as e:
            print(f"Error in analyze_data: {str(e)}")
            return {
                "title": "Analysis Error",
                "figures": [],
                "insights": ["Error generating analysis"],
                "analysis_plan": {},
                "data": df
            }
    
    def create_presentation(self, analysis_results: Dict[str, Any] = None) -> Presentation:
        """
        Create a PowerPoint presentation from the analysis results.
        
        Args:
            analysis_results: Optional analysis results, uses latest analysis if not provided
            
        Returns:
            PowerPoint presentation object
        """
        # Use provided results or fall back to latest results
        if analysis_results is None:
            if self.latest_analysis is None:
                raise ValueError("No analysis results available for presentation")
            
            # Regenerate the analysis if we have data but no figures
            analysis_results = self.analyze_data()
        
        # Create a new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = analysis_results["title"]
        subtitle = slide.placeholders[1]
        subtitle.text = f"Data Analysis Report\nGenerated on {pd.Timestamp.now().strftime('%Y-%m-%d')}"
        
        # Add insights slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Key Insights"
        
        # Add insights as bullet points
        content = slide.placeholders[1]
        tf = content.text_frame
        
        # Handle insights whether they're a string or list
        insights = analysis_results.get("insights", [])
        if isinstance(insights, str):
            # Split string into paragraphs
            insights = [p.strip() for p in insights.split('\n') if p.strip()]
        
        # Add each insight as a bullet point
        for insight in insights[:5]:  # Limit to 5 insights
            p = tf.add_paragraph()
            p.text = insight
            p.level = 0
        
        # Add visualization slides with analysis
        for i, fig in enumerate(analysis_results.get("figures", [])):
            # Save the figure to a bytes buffer
            buf = io.BytesIO()
            fig.savefig(buf, format='png', dpi=300, bbox_inches='tight')
            buf.seek(0)
            
            # Add a slide with a title, picture, and analysis
            slide_layout = prs.slide_layouts[5]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Add title from the analysis plan
            if "visualizations" in analysis_results.get("analysis_plan", {}):
                title = analysis_results["analysis_plan"]["visualizations"][i].get("description", f"Visualization {i+1}")
            else:
                title = f"Visualization {i+1}"
            slide.shapes.title.text = title
            
            # Add the figure (make it smaller to leave room for analysis)
            pic_left = Inches(0.5)
            pic_top = Inches(1.5)
            pic_width = Inches(6)  # Reduced width
            slide.shapes.add_picture(buf, pic_left, pic_top, width=pic_width)
            
            # Add analysis textbox
            analysis_left = Inches(7)  # Position to the right of the figure
            analysis_top = Inches(1.5)
            analysis_width = Inches(3)
            analysis_height = Inches(5)
            
            analysis_box = slide.shapes.add_textbox(analysis_left, analysis_top, analysis_width, analysis_height)
            tf = analysis_box.text_frame
            tf.word_wrap = True
            
            # Add analysis header
            p = tf.add_paragraph()
            p.text = "Analysis"
            p.font.bold = True
            p.font.size = Pt(14)
            
            # Add key observations
            if "visualizations" in analysis_results.get("analysis_plan", {}):
                viz_data = analysis_results["analysis_plan"]["visualizations"][i]
                if "analysis" in viz_data:
                    observations = viz_data["analysis"]
                    if isinstance(observations, list):
                        for obs in observations:
                            p = tf.add_paragraph()
                            p.text = f"• {obs}"
                            p.font.size = Pt(11)
                    else:
                        p = tf.add_paragraph()
                        p.text = observations
                        p.font.size = Pt(11)
            
            plt.close(fig)  # Close the figure to free memory
        
        # Add data summary slide if data is available
        if analysis_results.get("data") is not None and not analysis_results["data"].empty:
            df = analysis_results["data"]
            if len(df.columns) <= 10:  # Only add if not too many columns
                slide_layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Data Summary"
                
                # Create text box for data
                left = Inches(0.5)
                top = Inches(1.5)
                width = Inches(9)
                height = Inches(5)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.word_wrap = True
                
                # Add column headers
                p = tf.add_paragraph()
                p.text = " | ".join(str(col) for col in df.columns)
                p.font.bold = True
                
                # Add first few rows
                for _, row in df.head(10).iterrows():
                    p = tf.add_paragraph()
                    p.text = " | ".join(str(val) for val in row)
        
        # Add final slide
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        left = Inches(2.5)
        top = Inches(3)
        width = Inches(5)
        height = Inches(1)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        p = tf.add_paragraph()
        p.text = "Thank You"
        p.alignment = 1  # Center
        p.font.size = Pt(40)
        
        return prs
    
    def process_user_request(self, user_query: str) -> Dict[str, Any]:
        """
        End-to-end processing of a user query - from database query to presentation.
        
        Args:
            user_query: The user's natural language query
            
        Returns:
            Dictionary with query results, analysis, and path to presentation file
        """
        try:
            # Step 1: Query the database based on the user query
            data = self.query_database(user_query)
            
            # Step 2: Analyze the data and create visualizations
            analysis_results = self.analyze_data(data, user_query)
            
            # Step 3: Create a presentation
            presentation = self.create_presentation(analysis_results)
            
            # Save the presentation
            presentation_path = f"presentation_{pd.Timestamp.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            presentation.save(presentation_path)
            
            return {
                "success": True,
                "query": user_query,
                "data_shape": data.shape,
                "analysis_title": analysis_results["title"],
                "visualization_count": len(analysis_results["figures"]),
                "presentation_path": presentation_path
            }
            
        except Exception as e:
            return {
                "success": False,
                "error": str(e)
            }


# Example usage
if __name__ == "__main__":
    # Initialize with credentials from environment variables
    SUPABASE_URL = os.environ.get("SUPABASE_URL")
    SUPABASE_KEY = os.environ.get("SUPABASE_KEY")
    
    # Create the analysis system
    analyzer = DataAnalysisPresenter(SUPABASE_URL, SUPABASE_KEY)
    
    # Process a user query
    result = analyzer.process_user_request(
        "What insights can be drawn from the National AA (000s) and Coverage AA (000s) data in terms of audience engagement and growth over time?"
    )
    
    if result["success"]:
        print(f"Analysis complete! Presentation saved to: {result['presentation_path']}")
        print(f"Title: {result['analysis_title']}")
        print(f"Created {result['visualization_count']} visualizations")
    else:
        print(f"Error: {result['error']}")